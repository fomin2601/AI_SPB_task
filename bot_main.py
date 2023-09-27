# Fast refactor code from other participant

import telebot
import os
from telebot import types
from pptx import Presentation
from pptx.util import Inches
from get_ppt import PitchDeck
from utils import market, parsing

TOKEN = '6520533736:AAHtygH83eGv6AKZwoUYSiuT7-LFPt54ZI0'
bot = telebot.TeleBot(TOKEN)
answers = {}


@bot.message_handler(commands=['start'])
def start(message):
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    btn_yes = types.InlineKeyboardButton("да")
    btn_no = types.InlineKeyboardButton("нет")
    btn_start = types.InlineKeyboardButton("start")
    markup.add(btn_yes, btn_no)
    markup.add(btn_start)
    chat_id = message.chat.id
    bot.send_message(message.chat.id, 'Привет! Я помогу вам с созданием Pitch Deck. Готовы начать? (да/нет)', reply_markup=markup)


@bot.message_handler(func=lambda message: True)
def ask_question(message):
    chat_id = message.chat.id
    user_input = message.text.lower()
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    btn_yes = types.InlineKeyboardButton("да")
    btn_no = types.InlineKeyboardButton("нет")
    btn_start = types.InlineKeyboardButton("/start")
    markup.add(btn_start)

    if 'да' in user_input:
        bot.register_next_step_handler(message, scope_step)
    elif 'нет' in user_input:
        bot.send_message(message.chat.id, 'Хорошо, если у вас появятся вопросы, обращайтесь!')
    else:
        pass


def scope_step(message):
    scopes = get_ppt.codes
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    buttons = []

    for i, key in enumerate(scopes.keys()):
        buttons.append(types.InlineKeyboardButton(key))
        if (i + 1) % 4 == 0:
            markup.add(*buttons)
            buttons = []

    bot.send_message(message.chat.id, 'В какой области работает компания?', reply_markup=markup)
    bot.register_next_step_handler(message, company_name_step)


def company_name_step(message):
    answers['scope'] = message.text
    bot.send_message(message.chat.id, '1. Введите название компании', reply_markup=types.ReplyKeyboardRemove())
    bot.register_next_step_handler(message, short_description_step)


def short_description_step(message):
    answers['company_name'] = message.text
    bot.send_message(message.chat.id, '2. Введите краткое описание компании или введите "skip"')
    bot.register_next_step_handler(message, problem_definition_step)


def problem_definition_step(message):
    answers['short_description'] = message.text
    bot.send_message(message.chat.id, '3. Введите описание проблемы, которую решает проект, вашу целевую аудиторию')
    bot.register_next_step_handler(message, company_description_step)


def company_description_step(message):
    answers['problem_definition'] = message.text
    bot.send_message(message.chat.id, "4. Введите описание и ценностное предложение стартапа")
    bot.register_next_step_handler(message, solution_step)


def solution_step(message):
    answers['company_description'] = message.text
    bot.send_message(message.chat.id, "5. Введите предлагаемое решение поставленной проблемы")
    bot.register_next_step_handler(message, market_step)


def market_step(message):
    answers['solution'] = message.text
    bot.send_message(message.chat.id, "6. Введите рынок")
    bot.register_next_step_handler(message, competitors_step)


def competitors_step(message):
    answers['market'] = message.text
    bot.send_message(message.chat.id, "7. Введите ваших конкурентов")
    bot.register_next_step_handler(message, business_model_step)


def business_model_step(message):
    answers['competitors'] = message.text
    bot.send_message(message.chat.id, "8. Введите бизнес-модель")
    bot.register_next_step_handler(message, traction_and_finance_step)


def traction_and_finance_step(message):
    answers['business_model'] = message.text
    bot.send_message(message.chat.id, "9. Введите трекшн и финансы")
    bot.register_next_step_handler(message, team_step)


def team_step(message):
    answers['traction_and_finance'] = message.text
    bot.send_message(message.chat.id, "10. Введите вашу команду")
    bot.register_next_step_handler(message, investments_step)


def investments_step(message):
    answers['team'] = message.text
    bot.send_message(message.chat.id, "11. Введите, какой объём инвестиций требуется")
    bot.register_next_step_handler(message, roadmap_step)


def roadmap_step(message):
    answers['investments'] = message.text
    bot.send_message(message.chat.id, "12. Введите roadmap")
    bot.register_next_step_handler(message, contacts_step)


def contacts_step(message):
    answers['roadmap'] = message.text
    bot.send_message(message.chat.id, "13. Введите контактную информацию")
    bot.register_next_step_handler(message, pptx_creating_step)


def pptx_creating_step(message):
    answers['contacts'] = message.text
    print(answers)

    filename = f'{answers["company_name"]}.pptx'
    prs = Presentation()

    # Титульная страница
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title

    title.text = answers['company_name']

    if answers['short_description'] != 'skip':
        subtitle = slide.placeholders[1]
        subtitle.text = answers['short_description']

    # Описание проблемы
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Описание проблемы'
    subtitle.text = answers['problem_definition']

    # Описание стартапа
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Описание стартапа'
    subtitle.text = answers['company_description']

    # Описание проблемы
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Решение проблемы'
    subtitle.text = answers['solution']

    # Анализ рынка
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Анализ рынка'
    tam, sam, som = market.tam_sam_som(parsing.get_market_size(get_ppt.codes[answers['scope']]))
    subtitle.text = f'TAM: {tam}\nSAM: {sam}\nSOM: {som}'
    print('tam/sam/som calculated')

    # Конкуренты
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Конкуренты'
    subtitle.text = answers['competitors']

    # Бизнес-модель
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Бизнес-модель'
    subtitle.text = answers['business_model']

    # Трекшн и финансы
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Трекшн и финансы'
    subtitle.text = answers['traction_and_finance']

    # Команда
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Команда'
    subtitle.text = answers['team']

    # Инвестиции
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Инвестиционный раунд'
    subtitle.text = answers['investments']

    # Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Roadmap'
    subtitle.text = answers['roadmap']

    # Контакты
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = 'Контакты'
    subtitle.text = answers['contacts']

    prs.save(filename)
    # Отправка презентации пользователю
    with open(filename, 'rb') as f:
        bot.send_document(message.chat.id, f)

    # Удаление временного PDF-файла
    os.remove(filename)


# Запуск бота
bot.polling()
