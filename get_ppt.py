from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from utils import parsing

codes = {
    'Advertising & Marketing': ['731000'],
    'Aero & SpaceTech': ['511000', '512000'],
    'AgroTech': ['011000', '012000', '013000', '014000', '015000', '016000', '017000'],
    'AI': ['620900'],  # not actual
    'AR/VR': ['620900'],  # not actual
    'BeautyTech': ['960200'],
    'BigData': ['631000'],
    'Business Intelligence': ['620300'],
    'Business Software': ['620100'],
    'CleanTech': ['812200'],
    'ConstructionTech': ['411000'],
    'Consumer Goods & Services': ['471000', '472000', '475000', '476000', '477000', '478000'],
    'Cybersecurity': ['620200', '749090'],
    'E-commerce': ['582900', '479000'],
    'EdTech': ['854200', '853000', '852200'],
    'Energy': ['351000', '352000', '353000'],
    'FinTech': ['641000', '642000', '643000', '649000'],
    'FoodTech': ['561000', '562000', '563000'],
    'Gaming': ['582100'],
    'GreenTech': ['990000', '749050'],
    'Hardware': ['620300', '951000'],
    'HrTech': ['781000', '782000', '783000'],
    'IndustrialTech': ['411000', '412000', '439900'],
    'Legal & RegTech': ['691000', '692000'],
    'Mapping & Navigation': ['639900'],
    'Media & Entertainment': ['591000', '592000', '601000', '602000'],
    'MedTech': ['861000', '862000'],
    'Real Estate': ['681000', '682000', '683000'],
    'RetailTech': ['471000', '473000', '476000', '477000'],
    'SafetyTech': ['651000', '652000', '653000'],
    'SportTech': ['931000'],
    'Telecom & Communication': ['611000', '612000', '613000', '619000'],
    'Transport & Logistics': ['492000', '494000', '502000', '504000', '511000'],
    'Travel': ['960400'],
    'Web3': ['631000', '721000'],
    'WorkTech': ['732000', '701000', '702000'],
}

info = {
    'project_name': 'Unspot',
    'field': 'SportTech',
    'problem': '',
    'definition': """Система UnSpot – это программное обеспечение для бронирования,
                  позволяющее сотрудникам временно резервировать для себя любое незакрепленное 
                  рабочее место в офисе компании""",
    'solution': 'фывфыв',
    'market': 'фывфы',
    'competitors': 'фвыфв',
    'business_model': 'фвыфы',
    'traction_finance': 'фвфвы',
    'team_and_board': 'фывфыв',
    'investitions': 'фвыв',
    'roadmap': 'фывыфвфы',
    'contacts': '12345',
}


class PitchDeck:
    def __init__(self, style: str = 'default'):
        self.presentation = Presentation()
        self.style = style

    def create_slide_text_and_img(
            self,
            text: str = 'Text sample for your slide, you can replace it by yours',
            img_path: str = 'static/standard_img.png'
    ):
        """
        Add slide with textbox and image using pre-build pattern, inplace

        :param prs: presentation to modify
        :param text: text to place on slide
        :param img_path: img to place on slide
        """

        # Select blank slide layout
        blank_slide_layout = self.presentation.slide_layouts[6]

        # Attach blank slide to presentation
        slide = self.presentation.slides.add_slide(blank_slide_layout)

        # Margins for text
        left_text = Cm(1)
        top_text = Cm(5)
        width_text = Cm(9.54)
        height_text = Cm(4.75)

        # Add text
        text_box = slide.shapes.add_textbox(left=left_text, top=top_text, width=width_text, height=height_text)
        tf = text_box.text_frame
        tf.text = 'Test text of this slide'

        # Margins for image
        left_img = Cm(11)
        top_img = Cm(1)
        width_img = Cm(13)
        height_img = Cm(17)

        # Add image
        img = slide.shapes.add_picture(img_path, left=left_img, top=top_img, width=width_img, height=height_img)

    def save(self, path: str = 'test.pptx'):
        self.presentation.save(path)


def create_slide(prs, info, content, slide_type=1):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = content
    subtitle.text = info[content]


def build_pptx():
    prs = Presentation()
    print(prs.slide_layouts)
    title_slide_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = info['project_name']
    subtitle.text = "python-pptx was here!"

    for key in info:
        create_slide(prs, info, key, slide_type=1)

    prs.save('test.pptx')


prs = PitchDeck()

prs.create_slide_text_and_img()
prs.create_slide_text_and_img()

prs.save('test.pptx')
